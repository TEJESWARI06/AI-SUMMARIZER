import os
import string
import nltk
import spacy
import pdfplumber
from flask import Flask, render_template, request, send_file
from werkzeug.utils import secure_filename
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from pptx import Presentation
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.text_rank import TextRankSummarizer

# Initialize Flask app
app = Flask(__name__)

# Configure upload folder
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

# Download necessary NLTK resources
nltk.download("punkt")
nltk.download("stopwords")
nltk.download("wordnet")

# Load spaCy model once to avoid repeated loading
spacy.cli.download("en_core_web_sm")
nlp = spacy.load("en_core_web_sm")


def load_pptx(file_path):
    """Extract text from a PowerPoint file while preserving structure."""
    prs = Presentation(file_path)
    structured_content = {}

    for i, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else f"Slide {i}"
        subheadings = []
        bullet_points = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() and shape != slide.shapes.title:
                text = shape.text.strip()
                if text.isupper():  # Assuming uppercase text as subheading
                    subheadings.append(text)
                else:
                    bullet_points.append(text)

        structured_content[title] = {
            "subheadings": subheadings if subheadings else ["No subheadings available."],
            "bullet_points": bullet_points if bullet_points else ["No content available."],
        }
    return structured_content


def load_pdf(file_path):
    """Extract text from a PDF file."""
    structured_content = {"Document": {"subheadings": [], "bullet_points": []}}
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                structured_content["Document"]["bullet_points"].append(text)
    return structured_content


def preprocess_text(text):
    """Preprocess text: lowercasing, tokenization, stopword removal, and lemmatization."""
    text = text.lower().translate(str.maketrans("", "", string.punctuation))
    tokens = word_tokenize(text)
    stop_words = set(stopwords.words("english"))
    filtered_tokens = [word for word in tokens if word not in stop_words]

    # Use spaCy for lemmatization
    doc = nlp(" ".join(filtered_tokens))
    lemmatized_tokens = [token.lemma_ for token in doc]

    return " ".join(lemmatized_tokens)


def summarize_text(text, num_sentences=3):
    """Summarize text using TextRank algorithm and return structured bullet points."""
    if not text.strip() or text == "No content available.":
        return ["No content available."]

    parser = PlaintextParser.from_string(text, Tokenizer("english"))
    summarizer = TextRankSummarizer()
    summary = summarizer(parser.document, num_sentences)

    return [sentence._text for sentence in summary]


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return "No file part", 400

    file = request.files["file"]
    if file.filename == "":
        return "No selected file", 400

    if file and (file.filename.endswith(".pptx") or file.filename.endswith(".pdf")):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(file_path)

        try:
            if file.filename.endswith(".pptx"):
                structured_content = load_pptx(file_path)
            else:
                structured_content = load_pdf(file_path)

            summarized_content = {}
            for heading, content in structured_content.items():
                summarized_subheadings = [summarize_text(preprocess_text(sub)) for sub in content["subheadings"]]
                summarized_bullets = [summarize_text(preprocess_text(bullet)) for bullet in content["bullet_points"]]
                summarized_content[heading] = {
                    "subheadings": [item for sublist in summarized_subheadings for item in sublist],
                    "bullet_points": [item for sublist in summarized_bullets for item in sublist]
                }

            return render_template("result.html", summary=summarized_content, filename=filename)
        except Exception as e:
            return f"Error processing file: {str(e)}", 500

    return "Invalid file format. Please upload a .pptx or .pdf file.", 400


@app.route("/download/<filename>")
def download(filename):
    file_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    return send_file(file_path, as_attachment=True)


if __name__ == "__main__":
    app.run(debug=True)
